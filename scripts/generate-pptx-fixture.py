#!/usr/bin/env python3
"""Generate the real PPTX acceptance fixture used by the layout tests."""

from io import BytesIO
from pathlib import Path

from PIL import Image, ImageDraw
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt


ROOT = Path(__file__).resolve().parents[1]
OUTPUT = ROOT / "tests" / "fixtures" / "classroom-layouts.pptx"


def make_map_image():
    image = Image.new("RGB", (720, 420), "#e8f1e4")
    draw = ImageDraw.Draw(image)
    draw.rounded_rectangle((24, 24, 696, 396), radius=28, fill="#d9e9d2", outline="#276a58", width=8)
    draw.polygon([(130, 280), (245, 95), (360, 280)], fill="#8a6b45", outline="#5f452e")
    draw.line([(70, 330), (250, 300), (420, 345), (650, 290)], fill="#2f80ed", width=18)
    draw.ellipse((510, 70, 640, 200), outline="#276a58", width=7)
    draw.line((575, 70, 575, 200), fill="#276a58", width=4)
    draw.line((510, 135, 640, 135), fill="#276a58", width=4)
    stream = BytesIO()
    image.save(stream, format="PNG")
    stream.seek(0)
    return stream


def add_text_box(slide, text, left, top, width, height, size=24, color=(17, 24, 39), bold=False):
    shape = slide.shapes.add_textbox(left, top, width, height)
    paragraph = shape.text_frame.paragraphs[0]
    paragraph.alignment = PP_ALIGN.RIGHT
    run = paragraph.add_run()
    run.text = text
    run.font.name = "Arial"
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = RGBColor(*color)
    return shape


def build():
    presentation = Presentation()
    presentation.slide_width = Inches(13.333)
    presentation.slide_height = Inches(7.5)

    title_slide = presentation.slides.add_slide(presentation.slide_layouts[0])
    title_slide.background.fill.solid()
    title_slide.background.fill.fore_color.rgb = RGBColor(255, 247, 223)
    title_slide.shapes.title.text = "الحضارة المصرية القديمة"
    title_slide.placeholders[1].text = "نصوص وصور وأشكال وتخطيطات مختلفة"
    title_slide.shapes.title.text_frame.paragraphs[0].alignment = PP_ALIGN.RIGHT

    image_slide = presentation.slides.add_slide(presentation.slide_layouts[5])
    image_slide.background.fill.solid()
    image_slide.background.fill.fore_color.rgb = RGBColor(242, 247, 239)
    image_slide.shapes.title.text = "خريطة مظاهر السطح"
    image_slide.shapes.add_picture(make_map_image(), Inches(0.7), Inches(1.35), width=Inches(7.3), height=Inches(4.7))
    add_text_box(image_slide, "جبال — نهر — شبكة إحداثيات", Inches(8.2), Inches(2.0), Inches(4.4), Inches(1.5), 28, (23, 63, 55), True)
    add_text_box(image_slide, "صورة حقيقية داخل ملف PowerPoint", Inches(8.2), Inches(4.0), Inches(4.4), Inches(1.0), 20, (39, 106, 88))

    shapes_slide = presentation.slides.add_slide(presentation.slide_layouts[6])
    shapes_slide.background.fill.solid()
    shapes_slide.background.fill.fore_color.rgb = RGBColor(17, 24, 39)
    banner = shapes_slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.8), Inches(0.6), Inches(11.7), Inches(1.1))
    banner.fill.solid()
    banner.fill.fore_color.rgb = RGBColor(215, 173, 53)
    banner.line.color.rgb = RGBColor(255, 242, 191)
    banner.text = "أشكال موضوعة في أماكنها الأصلية"
    banner.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
    banner.text_frame.paragraphs[0].runs[0].font.size = Pt(30)
    banner.text_frame.paragraphs[0].runs[0].font.bold = True
    circle = shapes_slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(1.2), Inches(2.4), Inches(2.2), Inches(2.2))
    circle.fill.solid()
    circle.fill.fore_color.rgb = RGBColor(47, 128, 237)
    circle.line.color.rgb = RGBColor(255, 255, 255)
    triangle = shapes_slide.shapes.add_shape(MSO_SHAPE.ISOSCELES_TRIANGLE, Inches(5.1), Inches(2.25), Inches(2.6), Inches(2.5))
    triangle.fill.solid()
    triangle.fill.fore_color.rgb = RGBColor(138, 107, 69)
    triangle.rotation = 8
    arrow = shapes_slide.shapes.add_shape(MSO_SHAPE.RIGHT_ARROW, Inches(8.7), Inches(2.75), Inches(3.1), Inches(1.35))
    arrow.fill.solid()
    arrow.fill.fore_color.rgb = RGBColor(39, 106, 88)
    add_text_box(shapes_slide, "دائرة", Inches(1.55), Inches(5.0), Inches(1.5), Inches(0.6), 22, (255, 255, 255), True)
    add_text_box(shapes_slide, "مثلث", Inches(5.65), Inches(5.0), Inches(1.5), Inches(0.6), 22, (255, 255, 255), True)
    add_text_box(shapes_slide, "سهم", Inches(9.6), Inches(5.0), Inches(1.5), Inches(0.6), 22, (255, 255, 255), True)

    OUTPUT.parent.mkdir(parents=True, exist_ok=True)
    presentation.save(OUTPUT)
    print(OUTPUT)


if __name__ == "__main__":
    build()
